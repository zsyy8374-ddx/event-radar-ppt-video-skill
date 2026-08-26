#!/usr/bin/env python3
"""create_ppt · 重磅事件雷达模板 — 每次改写内容即可
⭐ 两种样式可选：
   样式A：深蓝科技风（0D1B2A） — 默认/事件雷达
   样式B：暗紫科技风（1A1A2E） — 紫苏雷达专属
   使用前取消/注释对应配色块即可
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# ═══════════════════════════════════════════════
# 🅰️ 样式A：深蓝科技风（默认/事件雷达/订单雷达）
# ═══════════════════════════════════════════════
# BG_DARK     = RGBColor(0x0D, 0x1B, 0x2A)
# BG_MID      = RGBColor(0x18, 0x2A, 0x3E)
# BAR_COLOR   = RGBColor(0x42, 0xA5, 0xF5)  # 亮蓝顶底栏
# ACCENT_GOLD = RGBColor(0xFF, 0xD5, 0x4F)
# ACCENT_BLUE = RGBColor(0x42, 0xA5, 0xF5)
# ACCENT_CYAN = RGBColor(0x80, 0xDE, 0xEA)
# LIGHT_GRAY  = RGBColor(0xD0, 0xDD, 0xEE)
# BOTTOM_BG   = RGBColor(0x08, 0x15, 0x25)
# LABEL_PREFIX = '重磅事件雷达 ⚡'
# COVER_TOP_LEFT = '⚡ 重磅事件雷达 | 日期'
# SITE_PREFIX = '⚡'

# ═══════════════════════════════════════════════
# 🅱️ 样式B：暗紫科技风（紫苏雷达专属）
# ═══════════════════════════════════════════════
# 主题含「紫苏雷达」「紫苏叶」「紫苏」时用此风格
BG_DARK     = RGBColor(0x1A, 0x1A, 0x2E)
BG_MID      = RGBColor(0x25, 0x25, 0x3F)
BAR_COLOR   = RGBColor(0x9B, 0x59, 0xB6)  # 紫苏紫顶底栏
ACCENT_GOLD = RGBColor(0xE8, 0xC5, 0x4A)  # 暖金
ACCENT_BLUE = RGBColor(0x5B, 0x9B, 0xD5)
ACCENT_CYAN = RGBColor(0x80, 0xDE, 0xEA)
ACCENT_ORANGE=RGBColor(0xFF, 0xAA, 0x44)
ACCENT_RED  = RGBColor(0xFF, 0x66, 0x66)
GREEN_OK    = RGBColor(0x66, 0xCC, 0x77)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xE8, 0xE8, 0xF0)
MID_GRAY    = RGBColor(0x99, 0xAA, 0xBB)
DIM_LABEL   = RGBColor(0x88, 0x99, 0xAA)
BOTTOM_BG   = RGBColor(0x12, 0x12, 0x22)
LABEL_PREFIX = '紫苏雷达 ·'
COVER_TOP_LEFT = '🍃 紫苏雷达 · 主题 | 2026.XX.XX'
SITE_PREFIX = '🍃'

# ═══════════════════════════════════════════════
# 通用函数（两种样式共用）
# ═══════════════════════════════════════════════

FS = 15        # 🚨 最小字号
BAR_TOP = 7.0  # 底部条位置

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_bg(s, c=BG_DARK):
    bg=s.background; fill=bg.fill; fill.solid(); fill.fore_color.rgb=c

def tb(s, l, t, w, h, txt, sz=FS, clr=WHITE, bold=False, align=PP_ALIGN.LEFT):
    box=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h))
    p=box.text_frame.paragraphs[0]; p.text=txt
    p.font.size=Pt(sz); p.font.color.rgb=clr; p.font.bold=bold
    p.font.name='Microsoft YaHei'; p.alignment=align
    box.text_frame.word_wrap=True; return box

def _run_para(p, segs, sz):
    for seg in segs:
        if isinstance(seg,str):
            if not seg.strip(): continue
            text,color,bold=seg,LIGHT_GRAY,False
        else:
            text=str(seg[0]); color=seg[1]; bold=seg[2] if len(seg)>2 else False
        if not text.strip(): continue
        run=p.add_run()
        run.text=text; run.font.size=Pt(sz); run.font.color.rgb=color
        run.font.bold=bold; run.font.name='Microsoft YaHei'

def paras(s, l, t, w, h, lines, sz=FS):
    """多段文字（每段可多色）"""
    box=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h))
    tf=box.text_frame; tf.word_wrap=True
    for idx,line in enumerate(lines):
        p=tf.paragraphs[0] if idx==0 else tf.add_paragraph()
        p.font.size=Pt(sz); p.font.name='Microsoft YaHei'; p.alignment=PP_ALIGN.LEFT
        p.space_after=Pt(4)
        _run_para(p, line if isinstance(line,list) else [line], sz)
    return box

def bar(s, l, t, w, h, c=BAR_COLOR):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb=c; sh.line.fill.background()

def card_bg(s, l, t, w, h):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb=BG_MID; sh.line.fill.background()

def title_line(s, title):
    bar(s, 0.6, 0.30, 0.10, 0.60)
    tb(s, 0.9, 0.25, 11, 0.55, title, 24, ACCENT_GOLD, True)

def bottom_bar(s, label=LABEL_PREFIX):
    bar(s, 0, BAR_TOP, 13.333, 0.50, BOTTOM_BG)
    bar(s, 0, BAR_TOP, 13.333, 0.04, BAR_COLOR)
    tb(s, 0.5, 7.08, 6, 0.3, f'{label}  |  Dengxian AI Research', 9, DIM_LABEL)

# ═══════════════════════════════════
# 【样式A】封面模板（深蓝科技风）
# ═══════════════════════════════════
# 取消注释即可用：
# s=prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
# bar(s, 0, 0, 13.333, 0.08, BAR_COLOR)
# bar(s, 0.6, 2.0, 2.5, 0.04, ACCENT_GOLD)
# tb(s, 0.6, 0.4, 11, 0.5, '⚡ 重磅事件雷达 | 日期', 18, ACCENT_CYAN, True)
# tb(s, 0.6, 2.3, 12, 2.0, '事件主题标题', 44, WHITE, True)
# tb(s, 0.6, 4.6, 11, 0.6, '核心摘要行', 20, ACCENT_ORANGE)
# tb(s, 0.6, 5.5, 11, 0.5, '关键数据行', FS, LIGHT_GRAY)
# tb(s, 0.6, 6.3, 11, 0.3, 'Dengxian AI Research', 12, MID_GRAY)
# bottom_bar(s)

# ═══════════════════════════════════
# 【样式B】封面模板（暗紫科技风 · 紫苏雷达）
# ═══════════════════════════════════
# 取消注释即可用：
# s=prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
# bar(s, 0, 0, 13.333, 0.50, BAR_COLOR)          # 紫色顶栏
# bar(s, 0.6, 2.0, 2.5, 0.04, ACCENT_GOLD)       # 金色横线
# tb(s, 0.6, 0.7, 11, 0.5, '🍃 紫苏雷达 · 主题 | 2026.XX.XX', 18, LIGHT_GRAY, True)
# tb(s, 0.6, 2.3, 12, 2.0, '主标题\\n副标题', 42, WHITE, True)
# tb(s, 0.6, 4.6, 11, 0.6, '核心摘要', 20, ACCENT_ORANGE)
# tb(s, 0.6, 5.4, 11, 0.5, '关键数据行', FS, LIGHT_GRAY)
# tb(s, 0.6, 5.9, 11, 0.4, '板块紫苏评估：X.X／10 ⭐⭐⭐', 18, ACCENT_GOLD, True)
# tb(s, 0.6, 6.3, 11, 0.3, '基于紫苏叶理论 | Dengxian AI Research', 12, MID_GRAY)
# bottom_bar(s)

# ═══════════════════════════════════
# 从这里开始写每个slide的内容
# ═══════════════════════════════════

# 示例：封面（样式B 紫苏雷达）
s=prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
bar(s, 0, 0, 13.333, 0.50, BAR_COLOR)
bar(s, 0.6, 2.0, 2.5, 0.04, ACCENT_GOLD)
tb(s, 0.6, 0.7, 11, 0.5, COVER_TOP_LEFT, 18, LIGHT_GRAY, True)
tb(s, 0.6, 2.3, 12, 2.0, '重磅事件雷达\n主题名称', 44, WHITE, True)
tb(s, 0.6, 4.6, 11, 0.6, '副标题/核心摘要', 20, ACCENT_ORANGE)
tb(s, 0.6, 5.5, 11, 0.5, '关键数据行', FS, LIGHT_GRAY)
tb(s, 0.6, 6.3, 11, 0.3, 'Dengxian AI Research', 12, MID_GRAY)
bottom_bar(s)

# 其他slide按需添加...

OUT = "/tmp/event_radar_pptv.pptx"
prs.save(OUT)
print(f"✅ PPT saved: {OUT} ({os.path.getsize(OUT)//1024}KB) — {len(prs.slides)} slides")
print(f"📐 最小字号: {FS}pt/{'暗紫科技风' if BG_DARK == RGBColor(0x1A, 0x1A, 0x2E) else '深蓝科技风'}/明亮文字")
